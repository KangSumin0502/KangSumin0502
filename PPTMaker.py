import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import wikipedia

# 슬라이드 하나에 텍스트 양
SLIDE_TEXT_LIMIT = 200 





def get_one_sentence(summary_text):
    if "." in summary_text:
        return summary_text.split(".")[0] + "."
    return summary_text



def clean_line_breaks(text):
    # \n이 2개 이상 연속되면 하나만 남김
    text = re.sub(r"\n{2,}", "\n", text)
    return text.strip()



def extract_keywords():
    with open("input.txt", "r", encoding="utf-8") as f:
        text = f.read()
    keywords = re.findall(r"\[(.*?)\]", text)
    # 중복 제거 + 순서 유지
    return list(dict.fromkeys(keywords))



def get_summary(keyword):
    try:
        wikipedia.set_lang("ko")
        return wikipedia.summary(keyword, sentences=4)
    except:
        return "해당 키워드에 대한 요약 정보를 찾을 수 없습니다."



def add_left_color_bar(slide, color=RGBColor(30, 50, 110)):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(0.4), Inches(7.5)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.color.rgb = color



def apply_title_style(slide, title_text):
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = title_text
    p = title.text_frame.paragraphs[0]
    p.font.name = "HY견고딕"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(30, 50, 110)
    p.alignment = PP_ALIGN.CENTER

    subtitle.text = ""



def create_toc_slide(prs, keywords):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_left_color_bar(slide)

    # 기본 텍스트 박스 제거
    for shape in list(slide.shapes):
        if shape.has_text_frame:
            slide.shapes._spTree.remove(shape._element)

    # 제목
    tx = slide.shapes.add_textbox(Inches(1), Inches(0.6), Inches(8), Inches(1))
    tf = tx.text_frame
    tf.text = "목차"
    p = tf.paragraphs[0]
    p.font.name = "HY견고딕"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(30, 50, 110)
    p.alignment = PP_ALIGN.LEFT

    # 목차 본문
    tx_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    tf2 = tx_box.text_frame
    tf2.clear()

    for i, kw in enumerate(keywords):
        p = tf2.add_paragraph()
        p.text = f"{i+1}. {kw}"
        p.font.name = "HY견고딕"
        p.font.size = Pt(32)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.alignment = PP_ALIGN.CENTER

    p = tf2.add_paragraph()
    p.text = f"{len(keywords)+1}. 전체 요약"
    p.font.name = "HY견고딕"
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.alignment = PP_ALIGN.CENTER

    p = tf2.add_paragraph()
    p.text = f"{len(keywords)+2}. QnA"
    p.font.name = "HY견고딕"
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.alignment = PP_ALIGN.CENTER



def create_keyword_content_slides(prs, keyword, summary):
    # 줄바꿈 정리부터 적용
    summary = clean_line_breaks(summary)

    # 문장으로 분리
    raw_sentences = summary.split(".")
    sentences = []
    for s in raw_sentences:
        s = s.strip()
        if s:
            sentences.append(s + ".")

    # 문장 길이 기반으로 슬라이드 분할
    slides_sentence_groups = []
    current_group = []
    current_len = 0

    for sent in sentences:
        sent_len = len(sent)

        if current_group and current_len + sent_len > SLIDE_TEXT_LIMIT:
            slides_sentence_groups.append(current_group)
            current_group = [sent]
            current_len = sent_len
        else:
            current_group.append(sent)
            current_len += sent_len

    if current_group:
        slides_sentence_groups.append(current_group)

    generated_slides = []

    # 이제 각 그룹마다 하나의 슬라이드 생성
    for group in slides_sentence_groups:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        add_left_color_bar(slide)

        # 제목
        title = slide.shapes.title
        title.text = keyword
        tp = title.text_frame.paragraphs[0]
        tp.font.name = "맑은 고딕"
        tp.font.size = Pt(36)
        tp.font.bold = True
        tp.font.color.rgb = RGBColor(30, 50, 110)
        tp.alignment = PP_ALIGN.LEFT

        # 본문
        body = slide.shapes.placeholders[1].text_frame
        body.clear()

        first = True
        for sent in group:
            p = body.add_paragraph()
            if first:
                # 첫 문장: 키워드 굵게 + 문장
                run_kw = p.add_run()
                run_kw.text = keyword + " - "
                run_kw.font.bold = True
                run_kw.font.size = Pt(22)
                run_kw.font.name = "맑은 고딕"

                run_tx = p.add_run()
                run_tx.text = sent
                run_tx.font.size = Pt(20)
                run_tx.font.name = "맑은 고딕"

                p.space_after = Pt(15)
                p.line_spacing = 1.3
                first = False
            else:
                # 나머지 문장: 그냥 문장만
                p.text = sent
                p.font.size = Pt(20)
                p.font.name = "맑은 고딕"
                p.space_after = Pt(15)
                p.line_spacing = 1.3

        generated_slides.append(slide)

    return generated_slides



def create_summary_slide(prs, keywords, summaries):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_left_color_bar(slide)

    title = slide.shapes.title
    title.text = "전체 요약"
    tp = title.text_frame.paragraphs[0]
    tp.font.name = "맑은 고딕"
    tp.font.size = Pt(36)
    tp.font.bold = True
    tp.font.color.rgb = RGBColor(30, 50, 110)
    tp.alignment = PP_ALIGN.LEFT

    body = slide.shapes.placeholders[1].text_frame
    body.clear()

    for kw in keywords:
        p = body.add_paragraph()

        run_kw = p.add_run()
        run_kw.text = kw
        run_kw.font.name = "맑은 고딕"
        run_kw.font.bold = True
        run_kw.font.size = Pt(20)

        run_tx = p.add_run()
        one_line = get_one_sentence(summaries[kw])
        run_tx.text = " : " + one_line
        run_tx.font.size = Pt(20)
        run_tx.font.color.rgb = RGBColor(60, 60, 60)

        p.space_after = Pt(10)
        p.line_spacing = 1.2



def create_qna_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tx = slide.shapes.add_textbox(Inches(3), Inches(3), Inches(4), Inches(1))
    tf = tx.text_frame
    tf.text = "QnA"

    p = tf.paragraphs[0]
    p.font.name = "HY견고딕"
    p.font.size = Pt(36)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER



def create_ppt(keywords, summaries):
    prs = Presentation()

    # 제목
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    apply_title_style(slide, keywords[0])

    # 목차
    create_toc_slide(prs, keywords)

    # 키워드 설명 슬라이드들
    for kw in keywords:
        create_keyword_content_slides(prs, kw, summaries[kw])

    # 요약
    create_summary_slide(prs, keywords, summaries)

    # QnA
    create_qna_slide(prs)

    prs.save("result.pptx")
    print("create result.pptx ")



def save_meta(keywords, summaries):
    with open("summury.txt", "w", encoding="utf-8") as f:
        for kw in keywords:
            f.write(f"[{kw}]\n요약: {summaries[kw]}\n\n")



def main():
    keywords = extract_keywords()
    if not keywords:
        print("input.txt에서 키워드를 찾을 수 없습니다.")
        return

    summaries = {kw: get_summary(kw) for kw in keywords}

    create_ppt(keywords, summaries)
    save_meta(keywords, summaries)


if __name__ == "__main__":
    main()
